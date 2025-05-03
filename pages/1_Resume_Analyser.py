import streamlit as st
import streamlit_authenticator as stauth
import PyPDF2
import re
import requests
from PyPDF2 import PdfReader
import json
import re
import time 
import io
import pandas as pd 
import base64
import uuid
import zipfile 
import google.generativeai as genai
import os
from pptx import Presentation
import yaml
from yaml.loader import SafeLoader
import ssl 
import certifi
from supabase import create_client, Client

from Login import get_authenticator
from genai_extract import extract_role_scores
import uuid
from supabase_client import supabase, BUCKET_NAME

if st.session_state.get("authentication_status") != True:
    st.warning("Please login to access this page.")
    st.stop()



# Inject logo and favicon
def add_favicon():
    #logo_path = ".static/cap.png"   # Path to your logo
    favicon_path = ".static/favicon.ico"  # Path to your favicon

    # Inject a favicon (using base64 encoding)
    with open(favicon_path, "rb") as f:
        favicon_data = f.read()
    favicon_base64 = base64.b64encode(favicon_data).decode()

    # Inject a logo at the top
    favicon_html = f"""
    <link rel="icon" type="image/x-icon" href="data:image/x-icon;base64,{favicon_base64}">
    """
    st.markdown(favicon_html, unsafe_allow_html=True)

add_favicon()


# Function to extract text from a PDF
def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text += page.extract_text()
    return text

# Function to clean and preprocess text
def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^a-zA-Z0-9\s]', '', text)
    text = text.lower()
    return text

# Function to call the Gemini LLM API
def ask_gemini(prompt):
    response = gemini_model.generate_content(prompt)
    return response.text


# --- Upload Function ---
def upload_to_supabase(uploaded_file):
    file_name = uploaded_file.name
    file_bytes = uploaded_file.read()
    file_path = f"uploads/{file_name}"
 
    # Simulate upsert: delete if it already exists
    try:
        supabase.storage.from_(BUCKET_NAME).remove([file_path])
    except Exception as e:
        pass  # Ignore errors (e.g. file not found)
 
    # Upload the file
    try:
        response = supabase.storage.from_(BUCKET_NAME).upload(
            path=file_path,
            file=file_bytes
        )
    except Exception as e:
        st.error(f"Upload failed: {e}")
        return None
 
    # Get public URL
    try:
        public_url = supabase.storage.from_(BUCKET_NAME).get_public_url(file_path)
        return public_url
    except Exception as e:
        st.error(f"Failed to generate public URL: {e}")
        return None


# Function to call the OpenAI API for matching JDs with CVs
def match_jds_with_cvs(jd_texts, cv_texts, weightage):

    scores = []
    results = []

    for jd_text in jd_texts:
        
        jd_scores = []
        jd_results = []

        for cv_text in cv_texts:
            prompt = f"""
            You are an expert in matching job descriptions (JDs) with resumes (CVs).
            Given the following job description and resume, assign separate relevance scores (between 0 and 1, 4 decimals only) for:
            - Technical skills
            - Functional skills
            - Managerial skills

            Based on:
            - Relevance of experience and skills for each competency area
            - Educational background and universities
            - Languages and other competencies
            - Profile summary

            Job Description: {jd_text}
            Resume: {cv_text}

            Use this JSON schema:
            ```json
            {{
                "name": "string (optional, default: 'Unknown')",
                "gender": "string (optional, default: 'male')",
                "nationality": "string (optional, default: 'Not provided')", 
                "languages": "string (optional, default: 'English')",
                "skills": "string (optional)",
                "education": "string (optional, extract only if present)",
                "technical_skills_score": "float (required, between 0 and 1)",
                "functional_skills_score": "float (required, between 0 and 1)",
                "managerial_skills_score": "float (required, between 0 and 1)",
                "relevance_score": "float (calculated as sum or weighted average of the above 3 scores)"
            }}
            ```
            """

            json_string_kpi = ask_gemini(prompt)

            try:
                json_dict_kpis = json.loads(json_string_kpi)
            except json.JSONDecodeError as e:
                json_dict_kpis = {"error": str(e)}

            # Extract individual scores
            tech_score = float(json_dict_kpis.get('technical_skills_score', 0))
            func_score = float(json_dict_kpis.get('functional_skills_score', 0))
            mgr_score = float(json_dict_kpis.get('managerial_skills_score', 0))

            # Final relevance score (sum or average)
            final_score = tech_score + func_score + mgr_score
            final_score = min(final_score, 1.0)  # Cap at 1.0 if needed

            # Apply weightage based on education
            for edu in weightage["criterion"]:
                if edu.lower() in cv_text.lower():
                    final_score *= weightage["criterion"][edu]
                    final_score = round(final_score, 4)

            json_dict_kpis["relevance_score"] = final_score

            jd_scores.append(final_score)
            jd_results.append(json_dict_kpis)

        scores.append(jd_scores)
        results.append(jd_results)
        
        # print("len(jd_texts):", len(jd_texts))
        # print("len(uploaded_jds):", len(uploaded_jds))
        # print("len(cv_texts):", len(cv_texts))
        # print("len(uploaded_cvs):", len(uploaded_cvs))
        # print("len(results):", len(results))

        # for idx, r in enumerate(results):
        #     print(f"len(results[{idx}]):", len(r))


        
        data = []

        # Make sure 'uploaded_jds' and 'uploaded_cvs' are available
        # They should be the list of uploaded files (with .name attribute)

        for i, jd_text in enumerate(jd_texts):
            if i >= len(results):
                continue  # No matching result for this JD

            # Get JD name safely
            jd_name = uploaded_jds.name if i < len(jd_texts) else f"JD_{i+1}"

            # Now, zip cv_texts and results[i] safely
            for j, (cv_text, result) in enumerate(zip(cv_texts, results[i])):
                # Get CV name safely
                cv_name = uploaded_cvs[j].name if j < len(uploaded_cvs) else f"CV_{j+1}"

                # safely work with 'result'
                tech_score = round(float(result.get('technical_skills_score', 0)), 4)
                func_score = round(float(result.get('functional_skills_score', 0)), 4)
                mang_score = round(float(result.get('managerial_skills_score', 0)), 4)
                total_score = round(float(result.get('relevance_score', 0)), 4)

                data.append([
                    jd_name,
                    cv_name,
                    result.get('name', 'Unknown'),
                    result.get('gender', 'male'),
                    result.get('nationality', 'Not provided'),
                    result.get('languages', 'English'),
                    result.get('skills', ''),
                    result.get('education', ''),
                    tech_score,
                    func_score,
                    mang_score,
                    total_score
                ])

        # Create DataFrame
        df = pd.DataFrame(data, columns=[
            "Job Description",
            "Resume",
            "Candidate Name",
            "Gender",
            "Nationality",
            "Languages",
            "Profile",
            "Education",
            "Technical Skills Score",
            "Functional Skills Score",
            "Managerial Skills Score",
            "Total Score"
        ])

        # Sort DataFrame
        df = df.sort_values(by=['Job Description', 'Total Score'], ascending=[True, False])


        return df

# Extract structured information from the provided CV text
def extract_information_from_cv(text):
    """Parses text to extract structured CV information using a pre-defined prompt."""
    prompt = f"""
        You are an expert in parsing CVs/resumes to extract structured information.

        From this text: '{text}', extract and organize the following information:

        - **Name:** Full name in CAPITAL LETTERS (if available).
        - **Gender:** Gender of the person (male or female).
        - **Nationality:** Mention nationality (default: "Not provided").
        - **Languages:** All spoken languages (default: "English").  
        - **Profile:** Summary highlighting career milestones in strictly 9 concise bullet points.  
        - **Experience:** Summarize work experience in strictly 12 concise bullet points covering main projects, each point should be between 25 to 45 words.  
        - **Competency:** Extract key competencies, tool names, programming languages and technologies in exactly 4 bullet points, combine where necessary.
        - **Education:** List degrees and certifications strictly up to 4 bullet points based only on available CV content.  

        Use this JSON schema:
        ```json
        {{
        "name": "string (optional, default: 'Unknown')",
        "gender": "string (optional, default: 'male')",
        "nationality": "string (optional, default: 'Not provided')", 
        "languages": "string (optional, default: 'English')",
        "profile": "string (optional)",
        "experience": "string (optional, 12 concise bullet points ensure each bullet point is between 25 to 45 words)",
        "competency": "string (optional, extract only if present, ensure each bullet point is max 20 words)",
        "education": "string (optional, extract only if present)"
        "location": string(required, default : "UAE")
        }}
        """
    
    json_string_kpi = ask_gemini(prompt)
    try:
        json_dict_kpis = json.loads(json_string_kpi)
    except json.JSONDecodeError as e:
        json_dict_kpis = {"error": str(e)}
    return json_dict_kpis


# Function to fill placeholders in PowerPoint slides
def fill_slide_template(prs, data):
    """
    Fill a PowerPoint presentation with data by replacing placeholders.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in data.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)
 
# Function to generate a PowerPoint file for each CV and save them in a zip file
def generate_individual_ppts(cv_data_list, template_path_male, template_path_female, output_folder):
    """
    Generate individual PowerPoint files for each CV and store them in a zip archive.
    """
    ppt_files = []
 
    for cv_data in cv_data_list:
        # Choose appropriate template based on gender
        if cv_data['gender'] == "male":
            prs = Presentation(template_path_male)
        else:
            prs = Presentation(template_path_female)
 
        # Fill the template with candidate data
        fill_slide_template(prs, cv_data)
 
        # Use the candidate's name for the filename; default to a generic name if missing
        candidate_name = cv_data.get("name", "Unnamed_Candidate").replace(" ", "_")
        ppt_filename = f"{candidate_name}.pptx"
        ppt_filepath = os.path.join(output_folder, ppt_filename)
 
        # Save the individual PowerPoint file
        prs.save(ppt_filepath)
        ppt_files.append(ppt_filepath)
        st.write(f"Generated PowerPoint for {candidate_name} saved at {ppt_filepath}")
 
    # Create a zip file to store all the PowerPoint files
    zip_filename = os.path.join(output_folder, "CVs.zip")
    with zipfile.ZipFile(zip_filename, "w") as zipf:
        for ppt_file in ppt_files:
            zipf.write(ppt_file, os.path.basename(ppt_file))
            os.remove(ppt_file)  # Remove individual ppt files after adding to zip
 
    return zip_filename


    
if "scores_df" not in st.session_state:
    st.session_state["scores_df"] = None

st.markdown("<br><br>",unsafe_allow_html = True)
st.write(f'Welcome **{st.session_state.get("name")}**')

# st.title("🚀 Intelligent CV and JD Matching Assistant")
# st.subheader("Empowering Talent Discovery with AI, One Match at a Time!")

st.markdown(
"""
<h1 style='text-align: center; color: #003366; font-size: 36px;'>
    🚀 Intelligent CV-JD Matching Assistant
</h1>
<h4 style='text-align: center; color: #666666; font-size: 18px;'>
    Empowering Talent Discovery with AI, One Match at a Time!
</h4>
""",
unsafe_allow_html=True)

# This line fixes spaces using pading - 1rem means less , 10 rem means more 
st.markdown(
    """
    <style>
        .block-container {
            padding-top: 2rem;
            padding-bottom: 2rem;
        }
    </style>
    """,
    unsafe_allow_html=True
)



st.write("Upload job descriptions and resumes to find the best matches.")


# Sidebar for API key and weights CSV upload
st.sidebar.image("logo.svg")
st.sidebar.title("Configuration")
with st.sidebar:
    user_api_key = st.text_input("Enter your Gemini API key:", type="password")
    uploaded_weights = st.file_uploader("Upload Weights CSV", type=["csv"])

    # Info message for user
    st.info(
        "ℹ️ If you do not upload a weights file, default weight values will be used automatically."
    )


    # --- Add Download Template Button ---
    sample_weights = pd.DataFrame({
        "description": ["harvard", "stanford", "mit", "oxford", "cambridge", "bachelor", "master", "phd"],
        "weight": [1.5, 1.4, 1.3, 1.2, 1.2, 1.1, 1.2, 1.3]
    })

    buffer = io.BytesIO()
    sample_weights.to_csv(buffer, index=False)
    buffer.seek(0)

    st.download_button(
        label="📄 Download Weights Template",
        data=buffer,
        file_name="weights_template.csv",
        mime="text/csv",
        help="Download a sample template to correctly format your weights CSV."
    )

    authenticator = get_authenticator()
    authenticator.logout() # logout button in sidebar 

    st.markdown(
    """
    <marquee behavior="scroll" direction="left" style="color:#0070AD; font-size:13px">
    Disclaimer™: An I&D Middle East tool.
    </marquee>
    """,
    unsafe_allow_html=True)


genai.configure(api_key=user_api_key)
gemini_model = genai.GenerativeModel('gemini-1.5-flash', generation_config={"response_mime_type": "application/json"})

def generate_role_scores_and_upload(uploaded_cv): 
    #cv_bytes = uploaded_cv.read()
    pdf_text = extract_text_from_pdf(uploaded_cv)
    cleaned_text = clean_text(pdf_text)
    #cv_text = cv_bytes.decode("utf-8", errors="ignore")  # or use pdfminer for actual text
    role_scores = extract_role_scores(cleaned_text)
 
    if role_scores:
        role_score_only = {k: v for k, v in role_scores.items() if isinstance(v, (int, float))}
        best_role = max(role_score_only, key=role_score_only.get)
        filename = uploaded_cv.name.replace(" ", "_")
        public_path = f"{best_role}/{uuid.uuid4()}_{filename}"
 
        # Check duplicate
        existing = supabase.table("cvs_table").select("*").eq("file_name", filename).execute()
        if existing.data:
            st.warning("This CV already exists.")
        else:
            # Upload file
            supabase.storage.from_(BUCKET_NAME).upload(public_path, uploaded_cv.read())
 
            # Get public URL
            url = supabase.storage.from_(BUCKET_NAME).get_public_url(public_path)
 
            # Save metadata
            supabase.table("cvs_table").insert({
                "name": role_scores.get("Name","Unknown"),
                "file_name": filename,
                "role_scores": role_score_only,
                "download_url": url
            }).execute()
            st.success("Uploaded and scored successfully!")
            
# def generate_role_scores_and_upload(uploaded_cv): 
#     try:
#         pdf_text = extract_text_from_pdf(uploaded_cv)
#         cleaned_text = clean_text(pdf_text)
#         role_scores = extract_role_scores(cleaned_text)
 
#         if role_scores:
#             best_role = max(role_scores, key=role_scores.get)
#             filename = uploaded_cv.name.replace(" ", "_")
#             public_path = f"{best_role}/{uuid.uuid4()}_{filename}"
 
#             # Check for duplicate in DB
#             try:
#                 existing = supabase.table("cvs_table").select("*").eq("file_name", filename).execute()
#                 if existing.data:
#                     st.warning(f"{filename} already exists in the database.")
#                     return
#             except Exception as e:
#                 st.error(f"❌ Error checking existing CV in Supabase: {e}")
#                 return
 
#             try:
#                 # Upload file to Supabase Storage
#                 supabase.storage.from_(BUCKET_NAME).upload(public_path, uploaded_cv.read())
 
#                 # Get public URL
#                 url = supabase.storage.from_(BUCKET_NAME).get_public_url(public_path)
 
#                 # Insert metadata into DB
#                 supabase.table("cvs_table").insert({
#                     "name": role_scores.get("Name", "Unknown"),
#                     "file_name": filename,
#                     "role_scores": role_scores,
#                     "download_url": url
#                 }).execute()
 
#                 st.success(f"✅ {filename} uploaded and scored successfully!")
 
#             except Exception as e:
#                 st.error(f"❌ Upload or DB insert failed for {filename}: {e}")
 
#     except Exception as e:
#         st.error(f"❌ Failed to process CV {uploaded_cv.name}: {e}")            


uploaded_jds = st.file_uploader("Upload Job Description", type=["pdf"], accept_multiple_files=False, 
                                key = st.session_state.get("jd_upload_key","jd_upload_default"))

uploaded_cvs = st.file_uploader("Upload Resumes", type=["pdf"], accept_multiple_files=True,
                                key = st.session_state.get("cv_upload_key","cv_upload_default"))

if "uploaded_jds" in st.session_state and uploaded_jds is None : 
    del st.session_state["uploaded_jds"]
    for key in ["scores_df"]:
        st.session_state.pop(key,None)

if "uploaded_cvs" in st.session_state and uploaded_cvs is None : 
    del st.session_state["uploaded_cvs"]
    for key in ["scores_df"]:
        st.session_state.pop(key,None)                                


col1 , col2, col3  = st.columns([2,1,.5])
with col1:
    submit_button = st.button("Submit", type = "primary")
with col3:
    reset_button = st.button("Reset All", type = "primary")
with col2:
    cg_cv_button = st.button("Generate Capgemini CVs", type = "primary")

st.markdown("<br>", unsafe_allow_html = True)       

if "submit_pressed" not in st.session_state:
    st.session_state.submit_pressed = False



if reset_button:
    for key in ["scores_df","uploaded_weights"]:
        st.session_state.pop(key,None)

    # clear file uploaders by changing keys 
        
    st.session_state["jd_upload_key"] = str(uuid.uuid4())
    st.session_state["cv_upload_key"] = str(uuid.uuid4())    
    st.rerun()       

if submit_button:
    # Check missing inputs and give proper warnings
    if not user_api_key:
        st.warning("🚨 Please enter your Gemini API key in the sidebar before proceeding.")
    elif not uploaded_jds:
        st.warning("📄 Please upload at least one Job Description (JD) PDF.")
    elif not uploaded_cvs:
        st.warning("📄 Please upload at least one Resume (CV) PDF.")
    

    if user_api_key and uploaded_jds and uploaded_cvs:

        st.session_state.submit_pressed =True # set to True when all are provided after submit button is pressed 

        with st.spinner("🕵️‍♂️🔍 Hang tight! I'm diving deep into the uploaded resumes & fishing out the best candidates for your job descriptions. 🐟✨ Stay tuned — greatness is being handpicked!"):
            # Adjust for single JD file
            jd_texts = []
            if uploaded_jds is not None:
                jd_text = clean_text(extract_text_from_pdf(uploaded_jds))
                jd_texts.append(jd_text)    

            cv_texts = [clean_text(extract_text_from_pdf(cv)) for cv in uploaded_cvs]

            # Load weightage from CSV file
            if uploaded_weights: 
                weightage_df = pd.read_csv(uploaded_weights)
                weightage = {"criterion": dict(zip(weightage_df['description'], weightage_df['weight']))}
            else:
                weightage = {
                    "criterion": {
                        "harvard": 1.5,
                        "stanford": 1.4,
                        "mit": 1.3,
                        "oxford": 1.2,
                        "cambridge": 1.2,
                        "bachelor": 1.1,
                        "master": 1.2,
                        "phd": 1.3
                    }
                }
            
            st.session_state['scores_df'] = match_jds_with_cvs(jd_texts, cv_texts, weightage)

# After Submit button is pressed 
if 'scores_df' in st.session_state and st.session_state['scores_df'] is not None:
    scores_df = st.session_state['scores_df'].round(4)

    # Make sure there's a Selected column
    if 'Selected' not in scores_df.columns:
        scores_df['Selected'] = False

    # Editable table
    edited_df = st.data_editor(
        scores_df,
        use_container_width=True,
        hide_index=True,
        column_config={
            "Selected": st.column_config.CheckboxColumn(required=False)
        }
    )

    # Get selected rows
    selected_rows_df = edited_df[edited_df['Selected'] == True]

    # Download either selected or all
    if not selected_rows_df.empty:
        download_df = selected_rows_df.drop(columns=['Selected'])
    else:
        download_df = edited_df.drop(columns=['Selected'])

    # Prepare CSV and Excel
    csv = download_df.to_csv(index=False).encode('utf-8')

    excel_buffer = io.BytesIO()
    with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
        download_df.to_excel(writer, index=False, sheet_name='Matching Scores')
    excel_data = excel_buffer.getvalue()

    # Prepare ZIP of selected CVs
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        # Collect CV filenames to include
        selected_cv_filenames = selected_rows_df['Resume'].tolist() if not selected_rows_df.empty else edited_df['Resume'].tolist()

        for uploaded_cv in uploaded_cvs:
            if uploaded_cv.name in selected_cv_filenames:
                # Add the file to the zip
                zip_file.writestr(uploaded_cv.name, uploaded_cv.getvalue())
    zip_buffer.seek(0)

    # do this only if submit button is pressed before and cg_cv_button is pressed 
    # if st.session_state.submit_pressed:
    #     if cg_cv_button:

    #         # Define folder path for templates
    #         templates_folder = "template"  # Path to the folder containing  templates
            
    #         # Check if the templates are available
    #         template_male_path = os.path.join(templates_folder, "cg_template_male.pptx")
    #         template_female_path = os.path.join(templates_folder, "cg_template_female.pptx")
            
    #         if not os.path.exists(template_male_path) or not os.path.exists(template_female_path):
    #             st.error("Template files not found in the specified folder. Please ensure both male_template.pptx and female_template.pptx are available.")
            
    #         else:
    #             with st.spinner("🎉 Baking the CVs in our awesome CG template!"):
    #                 # get results 
    #                 cv_data_list = []
    #                 selected_cv_filenames = selected_rows_df['Resume'].tolist() if not selected_rows_df.empty else edited_df['Resume'].tolist()
    #                 # Process each resume, extract information, and store the structured data
    #                 for uploaded_cv in uploaded_cvs:
    #                     # upload all cvs to supabase for the run  
    #                     if uploaded_cv is not None: # for all 
    #                         st.write("Generating role scores..")
    #                         generate_role_scores_and_upload(uploaded_cv)

                                
    #                 for uploaded_cv in uploaded_cvs:    
    #                     if uploaded_cv.name in selected_cv_filenames: # for selected or all 
    #                         pdf_text = extract_text_from_pdf(uploaded_cv)
    #                         cleaned_text = clean_text(pdf_text)
                            
    #                         llm_extraction = extract_information_from_cv(cleaned_text)
    #                         #print(str(resume) + " analysed !")
    #                         time.sleep(1)  # Throttle API requests to avoid hitting limits
    #                         cv_data_list.append(llm_extraction)
                
    #                 # Create an output folder in memory
    #                 output_folder = "/tmp/generated_ppts"  # Temporary folder for storing generated PPTs
    #                 os.makedirs(output_folder, exist_ok=True)
                
    #                 # Generate individual PowerPoint presentations and zip them
    #                 st.write("Generating CG ppt..")
    #                 zip_file_path = generate_individual_ppts(cv_data_list, template_male_path, template_female_path, output_folder)
    #                 # Provide download link for the zip file
    #                 col1, col2, col3 = st.columns([2, 2, 2])
    #                 with col2:
    #                     with open(zip_file_path, "rb") as f:
    #                         st.download_button(
    #                             label="Download CVs (CG Format)",
    #                             data=f,
    #                             file_name="CVs_CG_Format.zip",
    #                             mime="application/zip"
    #                         )
                                
    # else:
    #     st.warning("👉 Please press the **Submit** button first before generating CVs!")   

    if st.session_state.submit_pressed:
        if cg_cv_button:
     
            templates_folder = "template"
            template_male_path = os.path.join(templates_folder, "cg_template_male.pptx")
            template_female_path = os.path.join(templates_folder, "cg_template_female.pptx")
     
            if not os.path.exists(template_male_path) or not os.path.exists(template_female_path):
                st.error("Template files not found. Ensure both 'cg_template_male.pptx' and 'cg_template_female.pptx' exist.")
            else:
                with st.spinner("🎉 Baking the CVs in our awesome CG template!"):
                    cv_data_list = []
     
                    # Build the list of selected filenames
                    if not selected_rows_df.empty:
                        selected_cv_filenames = selected_rows_df['Resume'].tolist()
                    elif not edited_df.empty:
                        selected_cv_filenames = edited_df['Resume'].tolist()
                    else:
                        selected_cv_filenames = [cv.name for cv in uploaded_cvs]
     
                    st.write("Selected filenames for CG generation:", selected_cv_filenames)
     
                    # First pass: upload all CVs
                    for uploaded_cv in uploaded_cvs:
                        if uploaded_cv is not None:
                            st.write(f"🔍 Generating role scores for: {uploaded_cv.name}")
                            try:
                                generate_role_scores_and_upload(uploaded_cv)
                            except Exception as e:
                                st.error(f"❌ Failed to upload or score {uploaded_cv.name}: {str(e)}")
                                continue
     
                    # Second pass: generate CG slides only for selected CVs
                    for uploaded_cv in uploaded_cvs:
                        if uploaded_cv.name in selected_cv_filenames:
                            try:
                                st.write(f"🧠 Extracting information from: {uploaded_cv.name}")
                                pdf_text = extract_text_from_pdf(uploaded_cv)
                                cleaned_text = clean_text(pdf_text)
                                llm_extraction = extract_information_from_cv(cleaned_text)
     
                                if llm_extraction:
                                    cv_data_list.append(llm_extraction)
                                else:
                                    st.warning(f"No data extracted from {uploaded_cv.name}")
     
                                time.sleep(1)  # To avoid API limits
                            except Exception as e:
                                st.error(f"❌ Error processing {uploaded_cv.name}: {str(e)}")
     
                    if not cv_data_list:
                        st.error("🚫 No CVs matched or were extracted properly. Please check your input or selections.")
                    else:
                        output_folder = "/tmp/generated_ppts"
                        os.makedirs(output_folder, exist_ok=True)
     
                        st.write("📦 Generating CG PowerPoint files...")
                        try:
                            zip_file_path = generate_individual_ppts(cv_data_list, template_male_path, template_female_path, output_folder)
                            col1, col2, col3 = st.columns([2, 2, 2])
                            with col2:
                                with open(zip_file_path, "rb") as f:
                                    st.download_button(
                                        label="⬇️ Download CVs (CG Format)",
                                        data=f,
                                        file_name="CVs_CG_Format.zip",
                                        mime="application/zip"
                                    )
                        except Exception as e:
                            st.error(f"❌ Failed to generate PPTs: {str(e)}")
    else:
        st.warning("👉 Please press the **Submit** button first before generating CVs!")

    # Show download buttons
    col1, col2, col3 = st.columns([2, 2, 2])

    with col1:
        st.download_button(
            label="Download Scores (CSV)",
            data=csv,
            file_name='matching_scores.csv',
            mime='text/csv'
        )

    with col2:
        st.download_button(
            label="Download Scores (Excel)",
            data=excel_data,
            file_name='matching_scores.xlsx',
            mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )

    with col3:
        st.download_button(
            label="Download Selected CVs (ZIP)",
            data=zip_buffer,
            file_name='selected_cvs.zip',
            mime='application/zip'
        )

        

